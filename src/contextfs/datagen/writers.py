"""Format writers that materialise a :class:`FileSpec` onto disk.

Each writer takes a spec's ``content`` payload and produces a real file of the
corresponding type - a real PDF with a real page tree, a real OOXML document -
so that the Phase 5 extractors are exercised against genuine format complexity
rather than against text files with misleading extensions.

Determinism
-----------
PDFs are written with ReportLab's ``invariant`` mode, which suppresses the
creation-timestamp and document-ID fields that would otherwise differ between
runs, so PDF output is byte-identical across regenerations.

OOXML formats (DOCX/PPTX/XLSX) are ZIP containers whose entries carry
modification timestamps written by python-docx / python-pptx / openpyxl. Those
bytes therefore differ between runs even when the content is identical. The
corpus is consequently **content-deterministic, not byte-deterministic**: the
same paths, the same extracted text, and the same filesystem mtimes every time.
``scripts/verify_corpus.py`` checks determinism at the extracted-content level,
which is the level that actually matters for reproducing evaluation numbers.
"""

from __future__ import annotations

import os
from datetime import datetime
from pathlib import Path
from typing import Any

__all__ = ["write_file", "set_mtime", "WRITERS"]


def set_mtime(path: Path, when: datetime) -> None:
    """Set a file's access and modification time.

    Modification times are load-bearing in ContextFS: activity-session
    reconstruction (Phase 12) clusters on temporal proximity, and the
    metadata-consistency signal in date classification (Phase 10) compares a
    mentioned date against the document's mtime. A corpus with "now" as every
    timestamp would make both layers untestable.
    """
    stamp = when.timestamp()
    os.utime(path, (stamp, stamp))


# ---------------------------------------------------------------------------
# Plain text family
# ---------------------------------------------------------------------------


def write_text(path: Path, content: Any) -> None:
    """Write a UTF-8 text file (used for .txt, .md, and source code)."""
    path.write_text(str(content).strip() + "\n", encoding="utf-8", newline="\n")


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def write_pdf(path: Path, content: Any) -> None:
    """Write a multi-page PDF from ``[(heading, [paragraph, ...]), ...]``.

    Args:
        path: Destination file.
        content: Sequence of ``(heading, paragraphs)`` pairs. Each pair starts
            a new logical section; pages break automatically.
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer

    styles = getSampleStyleSheet()
    body = ParagraphStyle(
        "CfsBody", parent=styles["BodyText"], fontSize=10.5, leading=15, spaceAfter=6
    )
    heading = ParagraphStyle(
        "CfsHeading", parent=styles["Heading2"], fontSize=14, spaceBefore=10, spaceAfter=8
    )

    doc = SimpleDocTemplate(
        str(path),
        pagesize=A4,
        leftMargin=20 * mm,
        rightMargin=20 * mm,
        topMargin=20 * mm,
        bottomMargin=20 * mm,
        invariant=1,  # deterministic output: no creation date, fixed doc id
    )
    flow: list[Any] = []
    for index, (title, paragraphs) in enumerate(content):
        if index:
            flow.append(PageBreak())
        flow.append(Paragraph(title, heading))
        for para in paragraphs:
            flow.append(Paragraph(para, body))
            flow.append(Spacer(1, 2))
    doc.build(flow)


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def write_docx(path: Path, content: Any) -> None:
    """Write a Word document from ``[(style, text), ...]``.

    Args:
        path: Destination file.
        content: Sequence of ``(style, text)`` pairs where style is one of
            ``"h1"``, ``"h2"``, ``"p"``, or ``"bullet"``.
    """
    from docx import Document

    doc = Document()
    for style, text in content:
        if style == "h1":
            doc.add_heading(text, level=1)
        elif style == "h2":
            doc.add_heading(text, level=2)
        elif style == "bullet":
            doc.add_paragraph(text, style="List Bullet")
        else:
            doc.add_paragraph(text)
    doc.save(str(path))


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def write_pptx(path: Path, content: Any) -> None:
    """Write a slide deck from ``[(title, [bullet, ...]), ...]``."""
    from pptx import Presentation

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    for index, (title, bullets) in enumerate(content):
        layout = title_layout if index == 0 else bullet_layout
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        if index == 0:
            if bullets:
                slide.placeholders[1].text = "\n".join(bullets)
        else:
            frame = slide.placeholders[1].text_frame
            frame.text = bullets[0] if bullets else ""
            for bullet in bullets[1:]:
                para = frame.add_paragraph()
                para.text = bullet
                para.level = 1 if bullet.startswith("  ") else 0
    prs.save(str(path))


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------


def write_xlsx(path: Path, content: Any) -> None:
    """Write a workbook from ``[(sheet_name, [[cell, ...], ...]), ...]``.

    The first row of each sheet is treated as a header row and emboldened.
    Structured layout matters here: the Phase 10 classifier awards a
    ``structured_context`` signal to dates that appear inside tabular data, and
    a timetable is the canonical example of that.
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font

    wb = Workbook()
    wb.remove(wb.active)
    for sheet_name, rows in content:
        ws = wb.create_sheet(title=sheet_name[:31])
        for row in rows:
            ws.append(list(row))
        for cell in ws[1]:
            cell.font = Font(bold=True)
        for column_cells in ws.columns:
            width = max((len(str(c.value)) for c in column_cells if c.value is not None), default=8)
            ws.column_dimensions[column_cells[0].column_letter].width = min(max(width + 2, 10), 48)
    wb.save(str(path))


#: Dispatch table from :attr:`FileSpec.kind` to writer function.
WRITERS = {
    "txt": write_text,
    "md": write_text,
    "code": write_text,
    "pdf": write_pdf,
    "docx": write_docx,
    "pptx": write_pptx,
    "xlsx": write_xlsx,
}


def write_file(path: Path, kind: str, content: Any, mtime: datetime) -> None:
    """Materialise one corpus file and stamp its modification time.

    Args:
        path: Destination path; parent directories are created.
        kind: A key of :data:`WRITERS`.
        content: Format-specific payload (see the individual writers).
        mtime: Modification time to stamp onto the finished file.

    Raises:
        KeyError: If ``kind`` has no registered writer.
    """
    path.parent.mkdir(parents=True, exist_ok=True)
    WRITERS[kind](path, content)
    set_mtime(path, mtime)
