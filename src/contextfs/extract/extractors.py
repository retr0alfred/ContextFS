"""Per-format content extractors.

Every extractor has the same contract: take a path, return an
:class:`ExtractedDocument`, never raise. Format libraries are imported *inside*
the functions so that importing this module stays cheap - the CLI must not pay
for ``pypdf`` and ``python-pptx`` merely to print help.
"""

from __future__ import annotations

import csv
import io
from pathlib import Path

from contextfs.extract.base import ExtractedBlock, ExtractedDocument

__all__ = [
    "extract_text",
    "extract_code",
    "extract_pdf",
    "extract_docx",
    "extract_pptx",
    "extract_xlsx",
    "extract_csv",
    "extract_image",
    "read_text_file",
    "TEXT_ENCODINGS",
    "CODE_LANGUAGES",
]

#: Encodings tried in order when reading a text file. UTF-8 first because it is
#: correct for anything modern; cp1252 before latin-1 because on Windows it is
#: the far more likely legacy encoding and decodes smart quotes correctly.
TEXT_ENCODINGS = ("utf-8-sig", "utf-8", "cp1252", "latin-1")

#: Extension to language name, used to tag code documents.
CODE_LANGUAGES = {
    ".py": "Python",
    ".js": "JavaScript",
    ".ts": "TypeScript",
    ".java": "Java",
    ".c": "C",
    ".cpp": "C++",
    ".h": "C/C++ header",
    ".cs": "C#",
    ".go": "Go",
    ".rs": "Rust",
    ".rb": "Ruby",
    ".sh": "Shell",
    ".sql": "SQL",
    ".html": "HTML",
    ".css": "CSS",
    ".json": "JSON",
    ".yaml": "YAML",
    ".yml": "YAML",
}


def read_text_file(path: Path) -> tuple[str, str]:
    """Read a text file, trying several encodings.

    Args:
        path: File to read.

    Returns:
        ``(text, encoding_used)``. The final fallback decodes with replacement
        rather than failing, because a partially mojibaked document is still
        far more useful to retrieval than no document at all.
    """
    raw = path.read_bytes()
    for encoding in TEXT_ENCODINGS:
        try:
            return raw.decode(encoding), encoding
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace"), "utf-8 (with replacements)"


# ---------------------------------------------------------------------------
# Plain text and Markdown
# ---------------------------------------------------------------------------


def extract_text(path: Path, rel_path: str) -> ExtractedDocument:
    """Extract a .txt or .md file, splitting into blocks on blank lines.

    Markdown headings become heading blocks so the Phase 8 extractive
    summariser has something structural to prefer.
    """
    doc = ExtractedDocument(
        path=str(path), rel_path=rel_path, ext=path.suffix.lower(), extractor="text"
    )
    try:
        text, encoding = read_text_file(path)
    except OSError as exc:
        return ExtractedDocument.failed(path, rel_path, "text", f"read failed: {exc}")

    doc.meta["encoding"] = encoding
    doc.meta["line_count"] = text.count("\n") + 1

    blocks: list[ExtractedBlock] = []
    is_markdown = path.suffix.lower() in {".md", ".markdown"}
    in_table = False

    for index, chunk in enumerate(_split_paragraphs(text)):
        stripped = chunk.strip()
        if not stripped:
            continue
        # A Markdown pipe table is genuine tabular structure and must be
        # tagged as such - several corpus files put dates in one.
        lines = stripped.splitlines()
        table_like = is_markdown and sum(1 for ln in lines if ln.count("|") >= 2) >= 2
        heading = is_markdown and stripped.startswith("#")
        if table_like:
            in_table = True
        blocks.append(
            ExtractedBlock(
                index=index,
                kind="section" if heading else "paragraph",
                label=stripped.splitlines()[0][:60],
                text=stripped,
                is_tabular=table_like,
                is_heading=heading,
                row_count=len(lines) if table_like else 0,
            )
        )

    doc.blocks = blocks
    doc.meta["markdown_tables"] = in_table
    if not blocks:
        doc.warnings.append("file contained no non-blank text")
    return doc


def _split_paragraphs(text: str) -> list[str]:
    """Split text on blank lines, preserving internal line structure."""
    parts: list[str] = []
    current: list[str] = []
    for line in text.splitlines():
        if line.strip():
            current.append(line)
        elif current:
            parts.append("\n".join(current))
            current = []
    if current:
        parts.append("\n".join(current))
    return parts


# ---------------------------------------------------------------------------
# Source code
# ---------------------------------------------------------------------------


def extract_code(path: Path, rel_path: str) -> ExtractedDocument:
    """Extract a source file as a single block, tagged with its language.

    Code is deliberately *not* split into blocks per function. Retrieval here is
    about which file a user is looking for, not about locating a symbol; a
    per-function split would multiply the chunk count for no gain in re-finding
    accuracy while noticeably slowing embedding on this hardware.
    """
    doc = ExtractedDocument(
        path=str(path), rel_path=rel_path, ext=path.suffix.lower(), extractor="code"
    )
    try:
        text, encoding = read_text_file(path)
    except OSError as exc:
        return ExtractedDocument.failed(path, rel_path, "code", f"read failed: {exc}")

    language = CODE_LANGUAGES.get(path.suffix.lower(), "source")
    doc.meta.update(
        {
            "encoding": encoding,
            "language": language,
            "line_count": text.count("\n") + 1,
            "comment_lines": _count_comment_lines(text, path.suffix.lower()),
        }
    )
    if text.strip():
        doc.blocks = [
            ExtractedBlock(
                index=0,
                kind="section",
                label=f"{language} source: {path.name}",
                text=text.strip(),
            )
        ]
    else:
        doc.warnings.append("source file is empty")
    return doc


def _count_comment_lines(text: str, ext: str) -> int:
    """Count comment lines, best-effort, for metadata only."""
    markers = {"#"} if ext in {".py", ".sh", ".yaml", ".yml", ".rb"} else {"//", "--", "/*", "*"}
    return sum(
        1
        for line in text.splitlines()
        if any(line.strip().startswith(marker) for marker in markers)
    )


# ---------------------------------------------------------------------------
# CSV
# ---------------------------------------------------------------------------


def extract_csv(path: Path, rel_path: str) -> ExtractedDocument:
    """Extract a CSV file as one tabular block."""
    doc = ExtractedDocument(
        path=str(path), rel_path=rel_path, ext=path.suffix.lower(), extractor="csv"
    )
    try:
        text, encoding = read_text_file(path)
    except OSError as exc:
        return ExtractedDocument.failed(path, rel_path, "csv", f"read failed: {exc}")

    doc.meta["encoding"] = encoding
    try:
        rows = list(csv.reader(io.StringIO(text)))
    except csv.Error as exc:
        doc.warnings.append(f"csv parse degraded to plain text: {exc}")
        rows = [[line] for line in text.splitlines()]

    rendered = "\n".join(" | ".join(str(cell) for cell in row) for row in rows if any(row))
    doc.meta["row_count"] = len(rows)
    doc.meta["column_count"] = max((len(r) for r in rows), default=0)
    if rendered:
        doc.blocks = [
            ExtractedBlock(
                index=0,
                kind="sheet",
                label=f"csv: {path.name}",
                text=rendered,
                is_tabular=True,
                row_count=len(rows),
            )
        ]
    else:
        doc.warnings.append("csv contained no rows")
    return doc


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def extract_pdf(path: Path, rel_path: str) -> ExtractedDocument:
    """Extract a PDF page by page.

    A page that yields no text is recorded as a warning rather than dropped: on
    a real corpus that pattern means a scanned image, and the extraction report
    is where a user learns OCR would help. Silently returning fewer pages would
    hide it.
    """
    from pypdf import PdfReader
    from pypdf.errors import PdfReadError

    doc = ExtractedDocument(path=str(path), rel_path=rel_path, ext=".pdf", extractor="pdf")
    try:
        reader = PdfReader(str(path))
    except (PdfReadError, OSError, ValueError) as exc:
        return ExtractedDocument.failed(path, rel_path, "pdf", f"could not open: {exc}")

    if getattr(reader, "is_encrypted", False):
        try:
            reader.decrypt("")
        except Exception:  # noqa: BLE001 - any failure means we cannot read it
            return ExtractedDocument.failed(
                path, rel_path, "pdf", "encrypted PDF; no password available"
            )

    try:
        info = reader.metadata or {}
        doc.meta.update(
            {
                "page_count": len(reader.pages),
                "title": str(info.get("/Title", "") or ""),
                "author": str(info.get("/Author", "") or ""),
                "producer": str(info.get("/Producer", "") or ""),
            }
        )
    except Exception as exc:  # noqa: BLE001 - metadata is optional
        doc.warnings.append(f"metadata unreadable: {exc}")
        doc.meta["page_count"] = len(reader.pages)

    empty_pages = 0
    for index, page in enumerate(reader.pages):
        try:
            text = (page.extract_text() or "").strip()
        except Exception as exc:  # noqa: BLE001 - one bad page must not kill the file
            doc.warnings.append(f"page {index + 1} failed: {exc}")
            empty_pages += 1
            continue
        if not text:
            empty_pages += 1
            continue
        doc.blocks.append(
            ExtractedBlock(index=index, kind="page", label=f"page {index + 1}", text=text)
        )

    if empty_pages:
        doc.warnings.append(
            f"{empty_pages} of {doc.meta.get('page_count', '?')} page(s) yielded no text "
            "(likely scanned images; OCR is out of scope)"
        )
    if not doc.blocks:
        doc.ok = False
        doc.error = "no extractable text in any page"
    return doc


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def extract_docx(path: Path, rel_path: str) -> ExtractedDocument:
    """Extract a Word document, keeping tables tagged as tabular."""
    import docx
    from docx.opc.exceptions import PackageNotFoundError

    doc = ExtractedDocument(path=str(path), rel_path=rel_path, ext=".docx", extractor="docx")
    try:
        source = docx.Document(str(path))
    except (PackageNotFoundError, OSError, ValueError, KeyError) as exc:
        return ExtractedDocument.failed(path, rel_path, "docx", f"could not open: {exc}")

    index = 0
    for paragraph in source.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name if paragraph.style else "") or ""
        doc.blocks.append(
            ExtractedBlock(
                index=index,
                kind="section" if style.startswith("Heading") else "paragraph",
                label=style or "paragraph",
                text=text,
                is_heading=style.startswith("Heading") or style == "Title",
            )
        )
        index += 1

    for table_number, table in enumerate(source.tables):
        rows = [
            " | ".join(cell.text.strip() for cell in row.cells)
            for row in table.rows
            if any(cell.text.strip() for cell in row.cells)
        ]
        if not rows:
            continue
        doc.blocks.append(
            ExtractedBlock(
                index=index,
                kind="section",
                label=f"table {table_number + 1}",
                text="\n".join(rows),
                is_tabular=True,
                row_count=len(rows),
            )
        )
        index += 1

    try:
        properties = source.core_properties
        doc.meta.update(
            {
                "title": properties.title or "",
                "author": properties.author or "",
                "paragraph_count": len(source.paragraphs),
                "table_count": len(source.tables),
            }
        )
    except Exception as exc:  # noqa: BLE001 - properties are optional
        doc.warnings.append(f"core properties unreadable: {exc}")

    if not doc.blocks:
        doc.ok = False
        doc.error = "document contained no text or tables"
    return doc


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def extract_pptx(path: Path, rel_path: str) -> ExtractedDocument:
    """Extract a slide deck, one block per slide, plus speaker notes.

    Slide titles are marked as headings: on a deck, the title carries most of
    the retrievable meaning, and Phase 8's extractive summariser leans on that.
    """
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError

    doc = ExtractedDocument(path=str(path), rel_path=rel_path, ext=".pptx", extractor="pptx")
    try:
        deck = Presentation(str(path))
    except (PackageNotFoundError, OSError, ValueError, KeyError) as exc:
        return ExtractedDocument.failed(path, rel_path, "pptx", f"could not open: {exc}")

    index = 0
    empty_slides = 0
    for slide_number, slide in enumerate(deck.slides, start=1):
        parts: list[str] = []
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
                    if not title and shape == slide.shapes.title:
                        title = text.splitlines()[0]
            if getattr(shape, "has_table", False):
                rows = [
                    " | ".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows
                ]
                if rows:
                    doc.blocks.append(
                        ExtractedBlock(
                            index=index,
                            kind="slide",
                            label=f"slide {slide_number} table",
                            text="\n".join(rows),
                            is_tabular=True,
                            row_count=len(rows),
                        )
                    )
                    index += 1

        if parts:
            doc.blocks.append(
                ExtractedBlock(
                    index=index,
                    kind="slide",
                    label=f"slide {slide_number}: {title or 'untitled'}",
                    text="\n".join(parts),
                    is_heading=bool(title),
                )
            )
            index += 1
        else:
            empty_slides += 1

        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    doc.blocks.append(
                        ExtractedBlock(
                            index=index,
                            kind="notes",
                            label=f"slide {slide_number} notes",
                            text=notes,
                        )
                    )
                    index += 1
        except Exception as exc:  # noqa: BLE001 - notes are optional
            doc.warnings.append(f"slide {slide_number} notes unreadable: {exc}")

    doc.meta["slide_count"] = len(deck.slides)
    if empty_slides:
        doc.warnings.append(f"{empty_slides} slide(s) had no text (likely image-only)")
    if not doc.blocks:
        doc.ok = False
        doc.error = "deck contained no text"
    return doc


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------


def extract_xlsx(path: Path, rel_path: str) -> ExtractedDocument:
    """Extract a workbook, one tabular block per sheet.

    Opened with ``data_only=True`` so formula *results* are read rather than
    formula source. A cell reading ``=TODAY()+7`` tells retrieval nothing; the
    date it evaluates to is the fact a user remembers.

    Every sheet block is marked ``is_tabular=True``. This is the strongest input
    to Phase 10's structured-context signal, and it is why the corpus contains a
    spreadsheet of *incidental* dates as a control.
    """
    from openpyxl import load_workbook
    from openpyxl.utils.exceptions import InvalidFileException

    doc = ExtractedDocument(path=str(path), rel_path=rel_path, ext=".xlsx", extractor="xlsx")
    try:
        workbook = load_workbook(str(path), data_only=True, read_only=True)
    except (InvalidFileException, OSError, ValueError, KeyError) as exc:
        return ExtractedDocument.failed(path, rel_path, "xlsx", f"could not open: {exc}")

    try:
        sheet_names = list(workbook.sheetnames)
        doc.meta["sheet_count"] = len(sheet_names)
        doc.meta["sheet_names"] = sheet_names

        for index, name in enumerate(sheet_names):
            sheet = workbook[name]
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [_render_cell(value) for value in row]
                if any(cell for cell in cells):
                    rows.append(" | ".join(cells).rstrip(" |"))
            if not rows:
                doc.warnings.append(f"sheet {name!r} is empty")
                continue
            doc.blocks.append(
                ExtractedBlock(
                    index=index,
                    kind="sheet",
                    label=f"sheet: {name}",
                    text=f"{name}\n" + "\n".join(rows),
                    is_tabular=True,
                    row_count=len(rows),
                )
            )
    except Exception as exc:  # noqa: BLE001 - a malformed sheet must not abort the file
        doc.warnings.append(f"sheet iteration stopped early: {exc}")
    finally:
        workbook.close()

    if not doc.blocks:
        doc.ok = False
        doc.error = "workbook contained no non-empty sheets"
    return doc


def _render_cell(value: object) -> str:
    """Render a spreadsheet cell as text, preserving date readability."""
    if value is None:
        return ""
    if hasattr(value, "strftime"):
        # Render as ISO plus the day-month-year form the corpus persona uses,
        # so date detection sees a surface form it recognises either way.
        return value.strftime("%Y-%m-%d")
    return str(value).strip()


# ---------------------------------------------------------------------------
# Images (metadata only - no OCR, by design)
# ---------------------------------------------------------------------------


def extract_image(path: Path, rel_path: str) -> ExtractedDocument:
    """Record image metadata. No OCR - explicitly out of scope for this build.

    The filename is emitted as the block text because on a personal corpus it is
    frequently the only textual signal an image carries, and it is genuinely
    what a user remembers ("the screenshot of the timetable").
    """
    doc = ExtractedDocument(
        path=str(path), rel_path=rel_path, ext=path.suffix.lower(), extractor="image"
    )
    try:
        from PIL import Image

        with Image.open(path) as image:
            doc.meta.update(
                {
                    "format": image.format or "",
                    "mode": image.mode,
                    "width": image.width,
                    "height": image.height,
                }
            )
            exif = getattr(image, "_getexif", lambda: None)()
            if exif and 36867 in exif:  # DateTimeOriginal
                doc.meta["captured_at"] = str(exif[36867])
    except Exception as exc:  # noqa: BLE001 - Pillow raises many types
        doc.warnings.append(f"image metadata unreadable: {exc}")

    doc.blocks = [
        ExtractedBlock(
            index=0,
            kind="section",
            label="filename",
            text=path.stem.replace("_", " ").replace("-", " "),
        )
    ]
    doc.warnings.append("image content not extracted: OCR is out of scope for this build")
    return doc
